import json
import os
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import graphviz

# ---------------------------------------------------------
# Theme Color Palette (For presentation UI)
# ---------------------------------------------------------
COLOR_BG = RGBColor(248, 250, 252)         # Off-white / light slate canvas
COLOR_HEADER_BG = RGBColor(15, 23, 42)     # Deep Slate Navy banner
COLOR_ACCENT = RGBColor(14, 165, 233)      # Tech Cyan
COLOR_TEXT_DARK = RGBColor(30, 41, 59)     # Charcoal body text
COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # Slate grey subtext
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_CARD_BG = RGBColor(255, 255, 255)    # Card fill
COLOR_CARD_BORDER = RGBColor(226, 232, 240)
COLOR_CODE_BG = RGBColor(15, 23, 42)       # Terminal dark box

FONT_MAIN = "Arial"
FONT_CODE = "Courier New"

# ---------------------------------------------------------
# Helper Functions for Styling the presentation
# ---------------------------------------------------------
def apply_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header_banner(slide, title_text, subtitle_text=None):
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLOR_HEADER_BG
    banner.line.fill.background()

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(13.333), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = FONT_MAIN
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.name = FONT_MAIN
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_ACCENT

def add_card_box(slide, left, top, width, height, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
    else:
        card.line.fill.background()
    return card

# ---------------------------------------------------------
# Image & Graphviz Helpers (Sometimes an image will be generated on the title page, and also a wiring diagram will be generated on the assembly page)
# ---------------------------------------------------------
def fetch_component_image(query_name, save_path="component.png"):
    try:
        url = "https://en.wikipedia.org/w/api.php"
        params = {
            "action": "query",
            "generator": "search",
            "gsrsearch": query_name,
            "gsrlimit": 1,
            "prop": "pageimages",
            "pithumbsize": 600,
            "format": "json",
            "formatversion": 2
        }
        headers = {"User-Agent": "RoboticsMentorApp/1.0 (educational_project)"}
        res = requests.get(url, params=params, headers=headers, timeout=5)
        if res.status_code == 200:
            data = res.json()
            pages = data.get("query", {}).get("pages", [])
            if pages and "thumbnail" in pages[0]:
                img_url = pages[0]["thumbnail"]["source"]
                img_res = requests.get(img_url, headers=headers, timeout=5)
                if img_res.status_code == 200:
                    with open(save_path, "wb") as f:
                        f.write(img_res.content)
                    return save_path
    except Exception as e:
        print(f"Wikipedia fetch skipped: {e}")
    return None

def create_wiring_diagram_image(wiring_list, output_image_path="wiring_diagram.png"):
    dot = graphviz.Digraph(comment="Wiring Diagram", format="png")
    dot.attr(rankdir="LR", size="7,4", bgcolor="transparent")
    dot.attr("node", shape="box", style="filled,rounded", fillcolor="#0F172A", fontcolor="#FFFFFF", fontname="Arial", fontsize="10")
    dot.attr("edge", fontname="Arial", fontsize="9", color="#0EA5E9", fontcolor="#0284C7")

    for conn in wiring_list:
        # Pull dynamic source and target components
        f_comp = conn.get("from_component") or conn.get("component", "Component A")
        f_pin = conn.get("from_pin", "")
        t_comp = conn.get("to_component", "Microcontroller")
        t_pin = conn.get("to_pin", "")
        
        # Avoid self-referencing nodes if LLM makes a mistake
        if f_comp.strip().lower() != t_comp.strip().lower():
            label_text = f" {f_pin} ➔ {t_pin} "
            dot.edge(f_comp, t_comp, label=label_text)

    rendered_path = dot.render("temp_wiring", cleanup=True)
    os.rename(rendered_path, output_image_path)
    return output_image_path

# ---------------------------------------------------------
# Main Presentation Builder
# ---------------------------------------------------------
def create_deck(json_path="latest_guide.json", output_pptx="robotics_guide.pptx"):
    with open(json_path, "r") as f:
        data = json.load(f)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Pre-generate wiring diagram image if wiring exists
    wiring_list = data.get("wiring", [])
    diagram_png = None
    if wiring_list:
        try:
            diagram_png = create_wiring_diagram_image(wiring_list, "wiring_diagram.png")
        except Exception as e:
            print(f"Could not render Graphviz diagram: {e}")

    # =========================================================
    # SLIDE 1: Title
    # =========================================================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide1)

    top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.25))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_ACCENT
    top_bar.line.fill.background()

    add_card_box(slide1, Inches(0.8), Inches(0.8), Inches(7.5), Inches(5.9))

    tb = slide1.shapes.add_textbox(Inches(1.1), Inches(1.1), Inches(6.9), Inches(5.3))
    tf = tb.text_frame
    tf.word_wrap = True

    p_badge = tf.paragraphs[0]
    p_badge.text = f"DIFFICULTY: {data.get('difficulty_level', 'GENERAL').upper()}"
    p_badge.font.name = FONT_MAIN
    p_badge.font.size = Pt(11)
    p_badge.font.bold = True
    p_badge.font.color.rgb = COLOR_ACCENT

    p_title = tf.add_paragraph()
    p_title.text = data.get("project_title", "Robotics Project Guide")
    p_title.font.name = FONT_MAIN
    p_title.font.size = Pt(30)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_HEADER_BG

    p_sub = tf.add_paragraph()
    p_sub.text = f"\n{data.get('summary', '')}"
    p_sub.font.name = FONT_MAIN
    p_sub.font.size = Pt(13)
    p_sub.font.color.rgb = COLOR_TEXT_DARK

    hardware_list = data.get("hardware", [])
    if hardware_list:
        main_comp = hardware_list[0].get("name", "Arduino")
        img_file = fetch_component_image(main_comp, "title_comp.png")
        if img_file and os.path.exists(img_file):
            add_card_box(slide1, Inches(8.6), Inches(0.8), Inches(3.9), Inches(5.9))
            slide1.shapes.add_picture(img_file, Inches(8.9), Inches(1.6), width=Inches(3.3))

    # =========================================================
    # SLIDE 2: Bill of Materials
    # =========================================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide2)
    add_header_banner(slide2, "Bill of Materials", "Required hardware components & cost estimation")

    if hardware_list:
        add_card_box(slide2, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.3))

        rows = len(hardware_list) + 1
        table_shape = slide2.shapes.add_table(rows, 4, Inches(1.1), Inches(1.8), Inches(11.133), Inches(0.4 * rows))
        table = table_shape.table

        table.columns[0].width = Inches(4.5)
        table.columns[1].width = Inches(1.5)
        table.columns[2].width = Inches(2.0)
        table.columns[3].width = Inches(3.133)

        headers = ["Component", "Qty", "Est. Cost", "Notes"]
        for j, h_text in enumerate(headers):
            cell = table.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_HEADER_BG
            p = cell.text_frame.paragraphs[0]
            p.text = h_text
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = COLOR_WHITE

        for i, item in enumerate(hardware_list, start=1):
            row_bg = COLOR_WHITE if i % 2 == 1 else RGBColor(241, 245, 249)
            vals = [
                str(item.get("name", "")),
                str(item.get("quantity", 1)),
                f"${item.get('estimated_cost_usd', 0.0):.2f}",
                str(item.get("notes", "") or "—")
            ]
            for j, val in enumerate(vals):
                cell = table.cell(i, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_bg
                p = cell.text_frame.paragraphs[0]
                p.text = val
                p.font.size = Pt(11)
                p.font.color.rgb = COLOR_TEXT_DARK

# =========================================================
    # SLIDE 3: System Wiring Table
    # =========================================================
   
    slide3 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide3)
    add_header_banner(slide3, "Wiring & Circuit Pinout", "Complete schematic connections across components")

    if wiring_list:
        add_card_box(slide3, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.3))
        rows = len(wiring_list) + 1
        
        # 5 Columns: Source Component, From Pin, Target Component, To Pin, Wire/Purpose
        table_shape = slide3.shapes.add_table(rows, 5, Inches(1.0), Inches(1.7), Inches(11.333), Inches(0.38 * rows))
        table = table_shape.table

        table.columns[0].width = Inches(2.8)  # From Component
        table.columns[1].width = Inches(1.5)  # From Pin
        table.columns[2].width = Inches(2.8)  # To Component
        table.columns[3].width = Inches(1.5)  # To Pin
        table.columns[4].width = Inches(2.733)# Wire Purpose

        headers = ["From Component", "From Pin", "To Component", "To Pin", "Purpose / Signal"]
        for j, h_text in enumerate(headers):
            cell = table.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_HEADER_BG
            p = cell.text_frame.paragraphs[0]
            p.text = h_text
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_WHITE

        for i, conn in enumerate(wiring_list, start=1):
            row_bg = COLOR_WHITE if i % 2 == 1 else RGBColor(241, 245, 249)
            vals = [
                str(conn.get("from_component", conn.get("component", ""))),
                str(conn.get("from_pin", "")),
                str(conn.get("to_component", "Microcontroller")),
                str(conn.get("to_pin", "")),
                str(conn.get("wire_purpose", "Signal"))
            ]
            for j, val in enumerate(vals):
                cell = table.cell(i, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_bg
                p = cell.text_frame.paragraphs[0]
                p.text = val
                p.font.size = Pt(9.5)
                p.font.color.rgb = COLOR_TEXT_DARK

    # =========================================================
    # SLIDE 4: Assembly Steps & Wiring Diagram Visual
    # =========================================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide4)
    add_header_banner(slide4, "Assembly Instructions", "Step-by-step construction & circuit schematic")

    steps = data.get("steps", [])

    if diagram_png and os.path.exists(diagram_png):
        # 2-Column Layout: Left (Steps), Right (Diagram)
        add_card_box(slide4, Inches(0.8), Inches(1.5), Inches(6.0), Inches(5.3))
        tb4 = slide4.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.6), Inches(4.9))
        tf4 = tb4.text_frame
        tf4.word_wrap = True

        for idx, step in enumerate(steps):
            p = tf4.paragraphs[0] if idx == 0 else tf4.add_paragraph()
            p.text = f"STEP {step.get('step_number', idx+1)}: {step.get('title', '').upper()}"
            p.font.name = FONT_MAIN
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = COLOR_ACCENT

            p_desc = tf4.add_paragraph()
            p_desc.text = step.get("description", "")
            p_desc.font.name = FONT_MAIN
            p_desc.font.size = Pt(10.5)
            p_desc.font.color.rgb = COLOR_TEXT_DARK

            if step.get("safety_warning"):
                p_warn = tf4.add_paragraph()
                p_warn.text = f"⚠️ {step.get('safety_warning')}"
                p_warn.font.name = FONT_MAIN
                p_warn.font.size = Pt(9.5)
                p_warn.font.italic = True
                p_warn.font.color.rgb = RGBColor(225, 29, 72)

            tf4.add_paragraph().text = ""

        # Right side: Wiring Diagram Visual
        add_card_box(slide4, Inches(7.1), Inches(1.5), Inches(5.433), Inches(5.3))
        slide4.shapes.add_picture(diagram_png, Inches(7.3), Inches(1.8), width=Inches(5.0))

    else:
        # Full-width layout fallback if no diagram generated
        add_card_box(slide4, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.3))
        tb4 = slide4.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(11.133), Inches(4.9))
        tf4 = tb4.text_frame
        tf4.word_wrap = True

        for idx, step in enumerate(steps):
            p = tf4.paragraphs[0] if idx == 0 else tf4.add_paragraph()
            p.text = f"STEP {step.get('step_number', idx+1)}: {step.get('title', '').upper()}"
            p.font.name = FONT_MAIN
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = COLOR_ACCENT

            p_desc = tf4.add_paragraph()
            p_desc.text = step.get("description", "")
            p_desc.font.name = FONT_MAIN
            p_desc.font.size = Pt(11)
            p_desc.font.color.rgb = COLOR_TEXT_DARK

            if step.get("safety_warning"):
                p_warn = tf4.add_paragraph()
                p_warn.text = f"⚠️ Warning: {step.get('safety_warning')}"
                p_warn.font.name = FONT_MAIN
                p_warn.font.size = Pt(10)
                p_warn.font.italic = True
                p_warn.font.color.rgb = RGBColor(225, 29, 72)

            tf4.add_paragraph().text = ""

# =========================================================
    # SLIDE 5+: Source Code (Can be multiple slides if code is long enough)
    # =========================================================
    for code_file in data.get("code_files", []):
        raw_code = code_file.get("code", "// No code provided")
        lines = raw_code.split("\n")
        
        # Determine chunk size based on line count
        lines_per_slide = 22
        chunks = [lines[i:i + lines_per_slide] for i in range(0, len(lines), lines_per_slide)]
        total_parts = len(chunks)

        for part_idx, chunk in enumerate(chunks, start=1):
            slide_code = prs.slides.add_slide(blank_layout)
            apply_slide_background(slide_code)
            
            # Dynamic title naming for multi-part files
            part_suffix = f" (Part {part_idx}/{total_parts})" if total_parts > 1 else ""
            filename_header = f"Code: {code_file.get('filename', 'main.cpp')}{part_suffix}"
            
            add_header_banner(
                slide_code, 
                filename_header, 
                f"Language: {code_file.get('language', 'C++')}"
            )

            add_card_box(
                slide_code, 
                Inches(0.8), Inches(1.5), 
                Inches(11.733), Inches(5.3), 
                bg_color=COLOR_CODE_BG, 
                border_color=None
            )

            tb_code = slide_code.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(11.133), Inches(4.9))
            tf_code = tb_code.text_frame
            tf_code.word_wrap = False  # Avoids breaking long code lines awkwardly

            # Auto-scale font size if a chunk is dense
            font_size = Pt(9) if len(chunk) > 18 else Pt(10)

            code_text = "\n".join(chunk)
            p_code = tf_code.paragraphs[0]
            p_code.text = code_text
            p_code.font.name = FONT_CODE
            p_code.font.size = font_size
            p_code.font.color.rgb = RGBColor(226, 232, 240)

    prs.save(output_pptx)
    print(f"Theme presentation generated successfully: {output_pptx}")

if __name__ == "__main__":
    create_deck()
